import os
from flask import Flask, request, render_template, jsonify, send_file, make_response
from werkzeug.utils import secure_filename
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__, 
    template_folder='app/templates',
    static_folder='app/static'
)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'your-secret-key')

def create_word_doc(notes, filename):
    doc = Document()
    
    # Add title
    title = doc.add_paragraph('Speaker Notes')
    title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    title.runs[0].font.size = Pt(16)
    title.runs[0].font.bold = True
    doc.add_paragraph()  # Add a blank line
    
    # Add notes for each slide
    for i, note in enumerate(notes, 1):
        if note.strip():  # Only add non-empty notes
            # Add slide number as heading
            slide_header = doc.add_paragraph(f'Slide {i}')
            slide_header.style = 'Heading 2'
            
            # Add the notes
            doc.add_paragraph(note)
            doc.add_paragraph()  # Add a blank line between slides
    
    # Save to BytesIO object
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    
    return doc_io

def extract_notes(pptx_path):
    prs = Presentation(pptx_path)
    notes = []
    
    for slide in prs.slides:
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                notes_text += paragraph.text + "\n"
        notes.append(notes_text.strip())
    
    return notes

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not file.filename.endswith('.pptx'):
        return jsonify({'error': 'File must be a PowerPoint (.pptx) file'}), 400
    
    # Save the uploaded file
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    file.save(filepath)
    
    try:
        # Extract notes from PowerPoint
        notes = extract_notes(filepath)
        
        # Create Word document
        doc_io = create_word_doc(notes, filename)
        
        # Clean up the uploaded file
        os.remove(filepath)
        
        # Create response with both notes data and Word document
        response = make_response(jsonify({
            'success': True,
            'notes': notes,
            'doc_url': f'/download/{os.path.splitext(filename)[0]}_notes.docx'
        }))
        
        # Store the document in memory for download
        app.config['temp_docs'] = app.config.get('temp_docs', {})
        app.config['temp_docs'][f"{os.path.splitext(filename)[0]}_notes.docx"] = doc_io
        
        return response
    
    except Exception as e:
        # Clean up the uploaded file in case of error
        if os.path.exists(filepath):
            os.remove(filepath)
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>')
def download_file(filename):
    if 'temp_docs' not in app.config or filename not in app.config['temp_docs']:
        return jsonify({'error': 'File not found'}), 404
    
    doc_io = app.config['temp_docs'][filename]
    doc_io.seek(0)
    
    response = send_file(
        doc_io,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        as_attachment=True,
        download_name=filename
    )
    
    # Clean up the temporary document
    del app.config['temp_docs'][filename]
    
    return response

if __name__ == '__main__':
    app.run(debug=True) 