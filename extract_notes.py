from pptx import Presentation
from docx import Document
from docx.shared import Pt
import os

def extract_speaker_notes(pptx_path, output_path):
    """
    Extract speaker notes from a PowerPoint file and save them as a DOCX file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_path (str): Path where the DOCX file will be saved
    """
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Create a new Word document
    doc = Document()
    
    # Add title
    title = doc.add_heading('Speaker Notes', level=0)
    title.style = 'Title'
    doc.add_paragraph()  # Add a blank line
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        # Add slide number as header
        slide_header = doc.add_heading(f'Slide {i}', level=1)
        slide_header.style = 'Heading 1'
        
        # Get notes
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text.strip()
            if notes:
                # Add notes with proper formatting
                p = doc.add_paragraph()
                p.style = 'Normal'
                p.add_run(notes)
            else:
                p = doc.add_paragraph()
                p.style = 'Normal'
                p.add_run('No notes for this slide.')
        else:
            p = doc.add_paragraph()
            p.style = 'Normal'
            p.add_run('No notes for this slide.')
        
        # Add a blank line between slides
        doc.add_paragraph()
    
    # Save the document
    doc.save(output_path)

if __name__ == "__main__":
    # Get the current directory
    current_dir = os.path.dirname(os.path.abspath(__file__))
    
    # Define input and output paths
    pptx_file = os.path.join(current_dir, "Tabling Training 2025.pptx")
    output_file = os.path.join(current_dir, "speaker_notes.docx")
    
    # Extract notes
    extract_speaker_notes(pptx_file, output_file)
    print(f"Speaker notes have been extracted and saved to {output_file}") 