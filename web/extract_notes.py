from pptx import Presentation
from docx import Document
import os
import sys

def extract_speaker_notes(pptx_path, output_path):
    """
    Extract speaker notes from a PowerPoint file and save them as a DOCX file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_path (str): Path where the DOCX file will be saved
    """
    try:
        # Load the presentation
        prs = Presentation(pptx_path)
        
        # Create a new Word document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Speaker Notes', level=0)
        title.style = 'Title'
        doc.add_paragraph()  # Add a blank line
        
        # Process each slide
        for i, slide in enumerate(prs.slides, 1):
            # Add slide number as header
            slide_header = doc.add_heading(f'Slide {i}', level=1)
            slide_header.style = 'Heading 1'
            
            # Get notes
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()
                if notes:
                    # Add notes with proper formatting
                    p = doc.add_paragraph()
                    p.style = 'Normal'
                    p.add_run(notes)
                else:
                    p = doc.add_paragraph()
                    p.style = 'Normal'
                    p.add_run('No notes for this slide.')
            else:
                p = doc.add_paragraph()
                p.style = 'Normal'
                p.add_run('No notes for this slide.')
            
            # Add a blank line between slides
            doc.add_paragraph()
        
        # Save the document
        doc.save(output_path)
        print(f"Successfully processed {len(prs.slides)} slides")
        
    except Exception as e:
        print(f"Error processing PowerPoint file: {str(e)}", file=sys.stderr)
        raise

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_notes.py <input_pptx> <output_docx>", file=sys.stderr)
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    output_path = sys.argv[2]
    
    if not os.path.exists(pptx_path):
        print(f"Error: Input file '{pptx_path}' does not exist", file=sys.stderr)
        sys.exit(1)
    
    try:
        extract_speaker_notes(pptx_path, output_path)
    except Exception as e:
        print(f"Failed to process file: {str(e)}", file=sys.stderr)
        sys.exit(1) 